import argparse
import configparser
import hashlib
import io
import json
import os
import re
import sys
import textwrap
import google.generativeai as genai
import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tqdm import tqdm
from colorama import init, Fore, Style

# --- Constants ---
CACHE_DIR = ".cache"
BOX_WIDTH = 90  # Width of the output boxes

# --- Initialize Colorama for Windows compatibility ---
init(autoreset=True)

# --- Configuration & Caching ---
def load_config():
    """Loads settings from config.ini."""
    config = configparser.ConfigParser()
    config.read('config.ini')
    return config

def get_file_hash(file_path):
    """Calculates the MD5 hash of a file."""
    hasher = hashlib.md5()
    with open(file_path, 'rb') as f:
        buf = f.read()
        hasher.update(buf)
    return hasher.hexdigest()

def read_from_cache(cache_path):
    """Reads analysis results from a cache file."""
    with open(cache_path, 'r', encoding='utf-8') as f:
        return json.load(f)

def write_to_cache(cache_path, data):
    """Writes analysis results to a cache file."""
    os.makedirs(CACHE_DIR, exist_ok=True)
    with open(cache_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=4)

# --- Core Extraction Logic ---
def ocr_image(image_bytes):
    """Performs OCR on an image from bytes."""
    try:
        image = Image.open(io.BytesIO(image_bytes))
        return pytesseract.image_to_string(image)
    except Exception:
        return ""

def extract_content_from_pptx(file_path):
    """Extracts all textual content from a PowerPoint file."""
    if not os.path.exists(file_path):
        print(f"Error: The file '{file_path}' was not found.")
        sys.exit(1)

    prs = Presentation(file_path)
    slide_content = {}
    print("Extracting content from presentation...")
    for i, slide in enumerate(tqdm(prs.slides, desc="Processing Slides")):
        slide_number = i + 1
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                texts.append(shape.text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ocr_text = ocr_image(shape.image.blob)
                if ocr_text.strip():
                    texts.append(f"[Text from Image]:\n{ocr_text}")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            texts.append(f"[Speaker Notes]:\n{slide.notes_slide.notes_text_frame.text}")
        slide_content[slide_number] = "\n---\n".join(texts)
    return slide_content

# --- AI Analysis ---
def analyze_with_gemini(content_dict, api_key):
    """Analyzes content for inconsistencies using the Gemini API."""
    print("\nConnecting to Google AI for analysis...")
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-1.5-flash-latest')
    except Exception as e:
        return {"error": f"Failed to configure Gemini API: {e}"}

    prompt = f"""
    You are an expert business analyst and proofreader. Your task is to perform a deep analysis of the following PowerPoint presentation content.
    **Primary Objective:** Find all factual, numerical, and logical inconsistencies across the slides.
    **Secondary Objective:** If text is extracted from an image of a chart or table, try to interpret its structure and data. For example, if you see numbers in a table, verify if their sum matches any stated totals. Report any mathematical errors.
    **CRITICAL: You must format your entire response as a JSON object.** The root object should be a single key "issues" which contains a list of issue objects. Each issue object must have the following keys: "type", "conflict", and "evidence" (which is a list of strings).
    Example of a valid JSON response format:
    {{
      "issues": [
        {{
          "type": "Numerical Inconsistency",
          "conflict": "The total revenue stated on one slide does not match the sum of regional revenues on another.",
          "evidence": [
            "Slide 3 states: 'Total FY2024 Revenue: $10.2 Million'.",
            "Slide 8 chart shows regional revenues summing to $9.8 Million."
          ]
        }}
      ]
    }}
    If you find no issues at all, return an empty list: {{"issues": []}}
    Here is the presentation content:
    """
    full_text_prompt = "\n".join([f"--- Slide {slide_num} ---\n{text}" for slide_num, text in content_dict.items()])
    
    try:
        with tqdm(total=1, desc="Waiting for Gemini analysis") as pbar:
            response = model.generate_content(prompt + full_text_prompt)
            pbar.update(1)
        cleaned_text = response.text.strip().replace("```json", "").replace("```", "")
        return json.loads(cleaned_text)
    except json.JSONDecodeError:
        return {"error": "Failed to decode JSON from AI response. The model may have returned an invalid format."}
    except Exception as e:
        return {"error": f"An error occurred during the Gemini API call: {e}"}

# --- Colorful, Boxed Reporting ---
def generate_report(analysis_data):
    """Generates a visually enhanced report with colors and boxes."""
    if "error" in analysis_data:
        print(f"{Fore.RED}An error occurred:\n{analysis_data['error']}")
        return

    issues = analysis_data.get("issues", [])
    
    # --- Box Drawing Characters ---
    T_LEFT = '╔'
    T_RIGHT = '╗'
    B_LEFT = '╚'
    B_RIGHT = '╝'
    H_LINE = '═'
    V_LINE = '║'
    
    # --- Main Header ---
    print(f"\n{Fore.CYAN}{T_LEFT}{H_LINE * (BOX_WIDTH - 2)}{T_RIGHT}")
    title = " AI Inconsistency Analysis Report "
    print(f"{Fore.CYAN}{V_LINE}{Style.BRIGHT}{title.center(BOX_WIDTH - 2)}{Style.NORMAL}{V_LINE}")
    if issues:
        subtitle = f"Found {len(issues)} potential issue(s) to review."
        print(f"{Fore.CYAN}{V_LINE}{Fore.YELLOW}{subtitle.center(BOX_WIDTH - 2)}{Fore.CYAN}{V_LINE}")
    print(f"{Fore.CYAN}{B_LEFT}{H_LINE * (BOX_WIDTH - 2)}{B_RIGHT}\n")

    if not issues:
        print(f"{Fore.GREEN}✅ No inconsistencies were found in the presentation.\n")
        return

    # --- Issue Boxes ---
    for i, issue in enumerate(issues):
        issue_type = issue.get('type', 'N/A')
        conflict = issue.get('conflict', 'N/A')
        evidence = issue.get('evidence', [])

        # Issue Header
        header = f" ISSUE #{i+1} "
        print(f"{Fore.WHITE}{Style.BRIGHT}{header.center(BOX_WIDTH, '-')}")

        # Type
        type_line = f"{Fore.YELLOW}{Style.BRIGHT}TYPE:{Style.RESET_ALL} {issue_type}"
        print(type_line)

        # Conflict
        conflict_header = f"{Fore.RED}{Style.BRIGHT}CONFLICT:{Style.RESET_ALL} "
        wrapped_conflict = textwrap.wrap(conflict, width=BOX_WIDTH - len("CONFLICT: ") - 1)
        print(conflict_header + wrapped_conflict[0])
        for line in wrapped_conflict[1:]:
            print(" " * (len("CONFLICT: ") + 1) + line)

        # Evidence
        print(f"{Fore.CYAN}{Style.BRIGHT}EVIDENCE:{Style.RESET_ALL}")
        for ev in evidence:
            # Wrap each evidence line
            wrapped_evidence = textwrap.wrap(ev, width=BOX_WIDTH - 4) # 4 for "  - "
            print(f"  {Fore.WHITE}- {wrapped_evidence[0]}")
            for line in wrapped_evidence[1:]:
                print(f"    {line}")
        print("") # Add a newline for spacing

    print(f"{Fore.CYAN}{'=' * BOX_WIDTH}")


# --- Main Execution ---
def main():
    """Main function to orchestrate the analysis."""
    config = load_config()
    
    parser = argparse.ArgumentParser(description="Advanced PowerPoint Inconsistency Analyzer.", formatter_class=argparse.RawTextHelpFormatter)
    parser.add_argument("--file", help="Path to the .pptx file. Overrides 'default_file' in config.ini.")
    parser.add_argument("--output", help="Save the report to a file (e.g., report.txt, report.md).")
    parser.add_argument("--no-cache", action="store_true", help="Force re-analysis and ignore any cached results.")
    args = parser.parse_args()

    file_path = args.file or config.get('Settings', 'default_file', fallback=None)
    if not file_path:
        print("Error: No presentation file specified. Provide one with --file or in config.ini.")
        sys.exit(1)

    use_caching = not args.no_cache and config.getboolean('Options', 'caching', fallback=True)
    api_key = config.get('Settings', 'api_key', fallback=None)
    if not api_key or api_key == "YOUR_API_KEY_HERE":
        print("Error: API key not found in config.ini. Please add it.")
        sys.exit(1)

    analysis_result = None
    cache_path = ""
    if use_caching:
        try:
            file_hash = get_file_hash(file_path)
            cache_path = os.path.join(CACHE_DIR, f"{os.path.basename(file_path)}_{file_hash}.json")
            if os.path.exists(cache_path):
                print(f"Loading results from cache: {cache_path}")
                analysis_result = read_from_cache(cache_path)
        except FileNotFoundError:
            print(f"Error: The file '{file_path}' was not found.")
            sys.exit(1)
        except Exception as e:
            print(f"Could not read cache, proceeding with analysis. Error: {e}")

    if analysis_result is None:
        try:
            content = extract_content_from_pptx(file_path)
            if not content:
                print("Could not extract any content from the presentation.")
                return
            analysis_result = analyze_with_gemini(content, api_key)
            if use_caching and "error" not in analysis_result:
                write_to_cache(cache_path, analysis_result)
                print(f"Analysis complete. Results saved to cache: {cache_path}")
        except Exception as e:
            print(f"An unexpected error occurred during analysis: {e}")
            return

    # --- Generate and Display Report ---
    # Note: The colorful report is for terminal display only. 
    # File output will be plain text for now.
    if args.output:
        # We need a plain text version for file output
        # This part could be expanded to create formatted MD/HTML reports
        plain_report = generate_plain_text_report(analysis_result)
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(plain_report)
        print(f"\nReport successfully saved to {args.output}")
    else:
        generate_report(analysis_result)

def generate_plain_text_report(analysis_data):
    """Generates a plain text version of the report for file saving."""
    if "error" in analysis_data:
        return f"An error occurred:\n{analysis_data['error']}"
    issues = analysis_data.get("issues", [])
    if not issues:
        return "✅ No inconsistencies were found in the presentation."
    
    report_lines = [f"AI Inconsistency Analysis Report: Found {len(issues)} Issues\n{'='*80}"]
    for i, issue in enumerate(issues):
        report_lines.append(f"\n--- ISSUE #{i+1} ---\n")
        report_lines.append(f"TYPE: {issue.get('type', 'N/A')}")
        report_lines.append(f"CONFLICT: {issue.get('conflict', 'N/A')}")
        report_lines.append("EVIDENCE:")
        for ev in issue.get('evidence', []):
            report_lines.append(f"  - {ev}")
    return "\n".join(report_lines)


if __name__ == "__main__":
    main()
